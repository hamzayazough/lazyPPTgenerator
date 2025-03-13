import time
import pyautogui
from pptx import Presentation
from pptx.util import Inches
import pytesseract
import psutil
import os

def select_region(prompt):
    input(prompt)
    top_left = pyautogui.position()
    print(f"Top-left corner recorded at: {top_left}")
    input("Now, position your mouse at the BOTTOM-RIGHT corner of the region and press Enter...")
    bottom_right = pyautogui.position()
    print(f"Bottom-right corner recorded at: {bottom_right}")
    x = top_left.x
    y = top_left.y
    width = bottom_right.x - top_left.x
    height = bottom_right.y - top_left.y
    region = (x, y, width, height)
    print(f"Watching region: {region}")
    return region

def is_zoom_running():
    zoom_found = False
    for proc in psutil.process_iter(attrs=["name"]):
        name = proc.info.get("name")
        if name and "zoom" in name.lower():
            zoom_found = True
    return zoom_found

def get_slide_number(region):
    number_img = pyautogui.screenshot(region=region)
    text = pytesseract.image_to_string(number_img, config='--psm 7')
    try:
        return int(''.join(filter(str.isdigit, text)))
    except ValueError:
        return None

def capture_region(region):
    return pyautogui.screenshot(region=region)

def main():
    ppt_name = input("Enter the name you want to give to your ppt file (without extension): ")

    print("Please select the region containing the slide number.")
    slide_number_region = select_region(
        "Position your mouse over the TOP-LEFT corner of the slide number region and press Enter..."
    )
    
    print("Please select the region you want to capture as the slide.")
    slide_capture_region = select_region(
        "Position your mouse over the TOP-LEFT corner of the slide capture region and press Enter..."
    )
    
    ppt = Presentation()
    slide_layout = ppt.slide_layouts[6]
    
    last_slide_number = None
    print("Starting Zoom slide capture (press Ctrl+C to exit and save the presentation)")
    
    try:
        while True:
            if not is_zoom_running():
                print("Zoom is not running. Waiting for Zoom to start...")
                time.sleep(2)
                continue

            current_slide_number = get_slide_number(slide_number_region)
            
            if current_slide_number and current_slide_number != last_slide_number:
                last_slide_number = current_slide_number                
                slide_img = capture_region(slide_capture_region)
                print(f"Captured slide region for slide number: {current_slide_number}")
                image_filename = f"slide_{current_slide_number}.png"
                slide_img.save(image_filename)
                
                slide = ppt.slides.add_slide(slide_layout)
                left = top = Inches(1)
                slide.shapes.add_picture(image_filename, left, top, width=Inches(7.5))

                try:
                    os.remove(image_filename)
                except OSError as e:
                    print(f"Error removing file {image_filename}: {e}")
                
            time.sleep(2)
    except KeyboardInterrupt:
        print("You just ended the program. Saving the presentation...")
    
    downloads_dir = os.path.join(os.environ["USERPROFILE"], "Downloads")
    output_path = os.path.join(downloads_dir, f"{ppt_name}.pptx")
    ppt.save(output_path)
    print(f"Presentation saved as {output_path}")

if __name__ == "__main__":
    main()
